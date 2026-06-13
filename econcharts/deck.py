"""Assemble rendered charts into a PowerPoint deck — two charts per slide, side by
side. Called after a batch build so the user gets a ready deck named after the batch
(e.g. gallery.yaml -> gallery.pptx).

Each chart is placed at its EXACT export size (the named size's physical millimetres),
centered in its half of the slide — what-you-export is what-you-get on the slide. A
chart larger than its half-slide box is shrunk to fit (rare); transparent PNG
backgrounds blend onto the white slide.
"""

from __future__ import annotations

from pathlib import Path
from typing import Union

from pptx import Presentation
from pptx.util import Inches

# 16:9 widescreen slide; two side-by-side boxes with margins + a centre gap.
_SLIDE_W = Inches(13.333)
_SLIDE_H = Inches(7.5)
_MARGIN_X = Inches(0.4)
_MARGIN_Y = Inches(0.6)
_GAP = Inches(0.5)
_PER_SLIDE = 2
_MM_PER_IN = 25.4


def _emu_size(size_mm: tuple) -> tuple[int, int]:
    """(w_mm, h_mm) -> (width, height) in EMU."""
    w_mm, h_mm = size_mm
    return Inches(w_mm / _MM_PER_IN), Inches(h_mm / _MM_PER_IN)


def build_deck(items: list[tuple[Path, tuple]], out_path: Union[str, Path]) -> Path:
    """Build a .pptx, two charts per slide. `items` is (image_path, (w_mm, h_mm)) in
    order; each chart is placed at its real physical size. Returns the saved path."""
    prs = Presentation()
    prs.slide_width = _SLIDE_W
    prs.slide_height = _SLIDE_H
    blank = prs.slide_layouts[6]                 # the empty layout

    box_w = (_SLIDE_W - 2 * _MARGIN_X - _GAP) // 2
    box_h = _SLIDE_H - 2 * _MARGIN_Y
    lefts = [_MARGIN_X, _MARGIN_X + box_w + _GAP]

    for i in range(0, len(items), _PER_SLIDE):
        slide = prs.slides.add_slide(blank)
        for j, (img, size_name) in enumerate(items[i:i + _PER_SLIDE]):
            w, h = _emu_size(size_name)
            if w > box_w or h > box_h:           # safety: shrink to fit, keep aspect
                s = min(box_w / w, box_h / h)
                w, h = int(w * s), int(h * s)
            left = lefts[j] + (box_w - w) // 2
            top = _MARGIN_Y + (box_h - h) // 2
            slide.shapes.add_picture(str(img), left, top, width=w, height=h)

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return out_path
