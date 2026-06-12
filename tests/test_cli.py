"""CLI: `econcharts build` (fail-soft, exit codes) and `econcharts render`."""

from __future__ import annotations

import pytest

from econcharts.cli import main

_GOOD = """
output_dir: out
charts:
  - id: a
    title: A
    period: 2021Q1:2021Q4
    series: [{name: S, type: line, data: [1, 2, 3, 4]}]
  - id: b
    title: B
    period: 2021Q1:2021Q4
    series: [{name: S, type: bar, data: [1, 2, 3, 4]}]
"""


def _write(tmp_path, text, name="batch.yaml"):
    p = tmp_path / name
    p.write_text(text, encoding="utf-8")
    return p


def test_build_writes_figures_and_returns_zero(tmp_path, capsys):
    batch = _write(tmp_path, _GOOD)
    rc = main(["build", str(batch), "--force"])
    assert rc == 0
    out = tmp_path / "out"
    assert len(list(out.glob("a_*.png"))) == 1
    assert len(list(out.glob("b_*.png"))) == 1


def test_build_emits_pptx_deck_named_after_batch(tmp_path):
    from pptx import Presentation
    from pptx.util import Emu

    batch = _write(tmp_path, _GOOD)                  # 2 charts -> 1 slide, 2 pics
    rc = main(["build", str(batch), "--force"])
    assert rc == 0
    decks = list((tmp_path / "out").glob("*.pptx"))
    assert len(decks) == 1 and decks[0].stem == "batch"   # named after the yaml
    prs = Presentation(str(decks[0]))
    slides = list(prs.slides)
    assert len(slides) == 1
    pics = [sh for sh in slides[0].shapes if sh.shape_type == 13]  # PICTURE
    assert len(pics) == 2
    # placed at the chart's true physical size (default slides_half = 85 mm wide)
    assert abs(Emu(pics[0].width).inches - 85 / 25.4) < 0.02


def test_build_only_subset(tmp_path):
    batch = _write(tmp_path, _GOOD)
    rc = main(["build", str(batch), "--only", "a", "--force"])
    assert rc == 0
    out = tmp_path / "out"
    assert list(out.glob("a_*.png")) and not list(out.glob("b_*.png"))


def test_build_fail_soft_exit_code(tmp_path, capsys):
    text = _GOOD.replace("type: bar", "type: pie")  # chart b is now invalid
    batch = _write(tmp_path, text)
    rc = main(["build", str(batch), "--force"])
    assert rc == 1                                  # something failed
    assert list((tmp_path / "out").glob("a_*.png"))  # the good one still rendered
    assert "FAIL" in capsys.readouterr().err


def test_build_output_override(tmp_path):
    batch = _write(tmp_path, _GOOD)
    dest = tmp_path / "elsewhere"
    rc = main(["build", str(batch), "-o", str(dest), "--force"])
    assert rc == 0 and list(dest.glob("a_*.png"))


def test_render_single_spec(tmp_path):
    spec = _write(tmp_path, """
title: One
period: 2021Q1:2021Q4
series: [{name: S, type: line, data: [1, 2, 3, 4]}]
""", name="spec.yaml")
    out = tmp_path / "chart.png"
    rc = main(["render", str(spec), "-o", str(out)])
    assert rc == 0 and out.exists()


def test_render_bad_spec_returns_error(tmp_path):
    spec = _write(tmp_path, "title: X\nseries: []\n", name="bad.yaml")  # empty series
    rc = main(["render", str(spec), "-o", str(tmp_path / "x.png")])
    assert rc == 1
